"""Builds the Round 2 business proposal deck on the official AIC template.

    python3 docs/build_deck.py

The template is the one the challenge supplied and Round 1 used, so the deck
inherits its master, theme and fonts rather than inventing a look. Content comes
from docs/proposal.md; every figure in it traces to eval/results/report.md.

Regenerate rather than edit the .pptx by hand - the numbers change when the
benchmark is re-run, and a hand-edited deck goes stale silently.
"""
import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT.parent / "Breif_PPT" / "AIC_Talent-Brand_PPT-Template (1).pptx"
OUT = ROOT / "docs" / "ControlPlane_Business_Proposal.pptx"

# Pulled from the template's own palette so the deck stays on-brand.
INK = RGBColor(0x16, 0x20, 0x2B)
BODY = RGBColor(0x3A, 0x46, 0x52)
MUTED = RGBColor(0x6B, 0x77, 0x82)
ACCENT = RGBColor(0xA1, 0x00, 0xFF)      # Accenture purple
GOOD = RGBColor(0x1F, 0x7A, 0x5A)
BAD = RGBColor(0xA9, 0x39, 0x2C)
RULE = RGBColor(0xD8, 0xDD, 0xE3)
BAND = RGBColor(0xF4, 0xF1, 0xF8)


# --- slide plumbing -------------------------------------------------------

def layout(prs, master, name):
    for l in prs.slide_masters[master].slide_layouts:
        if l.name == name:
            return l
    raise KeyError(f"no layout {name!r} in master {master}")


def blank_copy(template):
    """The template with its slides removed, reloaded from disk.

    Dropping a slide relationship leaves the part in the package. Adding new
    slides afterwards hands them the partnames the dropped ones had, and the
    save writes both - which is how the first build produced two copies of one
    slide and lost another. Saving and reopening purges the orphans, so every
    slide added after this gets a clean name."""
    prs = Presentation(str(template))
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        prs.part.drop_rel(sldId.rId)
        lst.remove(sldId)
    tmp = OUT.with_name("._blank.pptx")
    prs.save(str(tmp))
    prs = Presentation(str(tmp))
    tmp.unlink()
    return prs


def cover_slide(prs, title, subtitle):
    s = prs.slides.add_slide(layout(prs, 0, "Cover: gradient"))
    for ph in s.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text_frame.text = title
            for r in ph.text_frame.paragraphs[0].runs:
                r.font.bold = True
        elif idx == 1:
            ph.text_frame.text = subtitle
    return s


def team_slide(prs):
    s = title_slide(prs, "Team Nexus", "IIT Jodhpur")
    members = [
        ("Akshat Jain", "M.Tech, Artificial Intelligence, 2027", "Round 2 implementation"),
        ("Aditya Pratap Singh", "M.Tech, Artificial Intelligence, 2027", ""),
        ("Arnesh Sanjeev Singh", "M.Tech, Computer Science, 2027", ""),
    ]
    top = 2.3
    for name, course, role in members:
        textbox(s, 0.61, top, 4.0, 0.9, [(name, 15, INK, True, 2)])
        textbox(s, 4.8, top + 0.04, 5.2, 0.9, [(course, 12, BODY, False, 0)])
        if role:
            textbox(s, 10.2, top + 0.04, 2.6, 0.9, [(role, 11.5, ACCENT, True, 0)])
        top += 0.92
    textbox(s, 0.61, 5.4, 12.1, 1.0, [
        ("Problem Track 1 — ControlPlane.ai · Accenture Innovation Challenge 2026",
         12, MUTED, False, 0)])
    return s


def salutation(prs, text="Thank you"):
    """A fresh closing slide on the template's own salutation layout. Cheaper
    than juggling slide indices to keep the original one in place, and it does
    not leave a stale title behind."""
    s = prs.slides.add_slide(layout(prs, 0, "Salutation: gradient"))
    if s.shapes.title is not None:
        s.shapes.title.text = text
    return s


def title_slide(prs, title, sub=None):
    s = prs.slides.add_slide(layout(prs, 4, "Content: title only"))
    s.shapes.title.text = title
    for p in s.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size, r.font.bold, r.font.color.rgb = Pt(26), True, INK
    if sub:
        box = s.shapes.add_textbox(Inches(0.61), Inches(1.16), Inches(12.1), Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = sub
        r = tf.paragraphs[0].runs[0]
        r.font.size, r.font.color.rgb = Pt(12.5), MUTED
    return s


def textbox(slide, left, top, width, height, blocks):
    """blocks: (text, size, colour, bold, space_after) tuples."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, colour, bold, after) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(after)
        for r in p.runs:
            r.font.size, r.font.color.rgb, r.font.bold = Pt(size), colour, bold
    return box


def table(slide, left, top, width, rows, col_widths, *, head=True,
          size=11, head_size=9.5, emphasis=None, colours=None):
    """A plain table. python-pptx's built-in styles fight the template, so the
    banding and rules are set explicitly and kept quiet."""
    n_rows, n_cols = len(rows), len(rows[0])
    height = Inches(0.32 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   Inches(width), height)
    tbl = shape.table
    tbl.first_row = head
    tbl.horz_banding = False
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.30 if r else 0.34)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            highlight = emphasis is not None and r == emphasis
            cell.fill.fore_color.rgb = (BAND if (r == 0 and head) or highlight
                                        else RGBColor(0xFF, 0xFF, 0xFF))
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if c and str(val)[:1].isdigit() else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(head_size if r == 0 and head else size)
                run.font.bold = bool((r == 0 and head) or highlight)
                run.font.color.rgb = (
                    MUTED if r == 0 and head else
                    (colours or {}).get((r, c), INK if highlight else BODY))
    return shape


def note(slide, text, top=6.62):
    box = slide.shapes.add_textbox(Inches(0.61), Inches(top), Inches(12.1), Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    r = tf.paragraphs[0].runs[0]
    r.font.size, r.font.color.rgb, r.font.italic = Pt(9.5), MUTED, True


# --- the deck -------------------------------------------------------------

def build():
    if not TEMPLATE.exists():
        sys.exit(f"template not found: {TEMPLATE}")
    prs = blank_copy(TEMPLATE)

    cover_slide(prs, "ControlPlane",
                "Business proposal · Round 2 · Problem Track 1\n"
                "Team Nexus, IIT Jodhpur")
    team_slide(prs)

    # 1 · the problem -----------------------------------------------------
    s = title_slide(prs, "The problem",
                    "One wrong sentence, and the loop that should catch it closes two weeks later.")
    textbox(s, 0.61, 1.95, 6.0, 4.4, [
        ("The policy document says room rent is capped at 1 percent.", 14, INK, True, 6),
        ("The assistant told the customer 2 percent.", 14, BAD, True, 14),
        ("Nothing in that answer looks wrong. It is worded exactly like a correct "
         "one. The customer acts on it. The enterprise finds out from a complaint.", 12, BODY, False, 14),
        ("The model answered in two seconds. The review reached the enterprise in "
         "two weeks. That gap is the problem.", 12, INK, True, 0),
    ])
    textbox(s, 7.0, 1.95, 5.7, 4.4, [
        ("Three things make it harder than it looks", 11.5, ACCENT, True, 10),
        ("The risks overlap. A fabricated detail about a named customer is a "
         "hallucination and a privacy event at once. One label per response loses that.", 11, BODY, False, 9),
        ("There is no ground truth at runtime. The knowledge gap that causes a "
         "hallucination is the same one that makes it hard to check.", 11, BODY, False, 9),
        ("Over- and under-flagging trade against each other. This cannot be solved, "
         "only priced — a missed hallucination in a regulated decision and a false "
         "alarm on an internal chatbot are not the same event.", 11, BODY, False, 0),
    ])

    # 2 · the cascade -----------------------------------------------------
    s = title_slide(prs, "What we built",
                    "A verification layer between the application and any model API. "
                    "Existing clients change one line: the base URL.")
    table(s, 0.61, 1.95, 12.1, [
        ["Tier", "What it does", "Cost", "Measured latency"],
        ["0", "Deterministic rules: PII patterns, access control, schema", "free", "under 1 ms"],
        ["1", "NLI entailment against each retrieved source, toxicity, bias heuristic",
         "GPU seconds", "87–190 ms p95, by policy"],
        ["2", "LLM judge, only for claims inside the policy's uncertainty band",
         "tokens", "1.3–3.1 s"],
    ], [0.7, 7.2, 1.7, 2.5])
    textbox(s, 0.61, 3.75, 12.1, 2.6, [
        ("Tier 2 fires on 2.8% of responses. That is not a target we set — it is what "
         "the uncertainty band produces, and the benchmark sweeps the band to show "
         "what every other setting would cost.", 12, BODY, False, 12),
        ("The action router does not compare a score to a threshold.", 12, INK, True, 4),
        ("It weighs P(wrong) × cost_of_being_wrong against cost_of_human_review and "
         "picks the cheaper mistake. Every category is evaluated, so a response "
         "flagged for both privacy and grounding takes the more severe action and "
         "records both reasons.", 12, BODY, False, 12),
        ("Every verdict, and every human override, is written to an append-only "
         "hash-chained log. A record cannot be altered afterwards without breaking "
         "the chain, and tampering is detected by test.", 12, BODY, False, 0),
    ])

    # 3 · the policy layer ------------------------------------------------
    s = title_slide(prs, "One response, three policies, three verdicts",
                    "Nothing in the code knows what a use case is.")
    table(s, 0.61, 2.0, 12.1, [
        ["", "customer support", "internal copilot", "decision support"],
        ["being wrong costs", "Rs 400", "Rs 60", "Rs 50,000"],
        ["hallucinated number", "regenerate", "annotate", "escalate"],
        ["PII leak", "block", "annotate", "block"],
        ["access-control breach", "block", "block", "block"],
        ["biased decision", "escalate", "annotate", "escalate"],
    ], [3.1, 3.0, 3.0, 3.0], emphasis=2)
    textbox(s, 0.61, 4.35, 12.1, 2.0, [
        ("Risk appetite, latency budget, uncertainty band, thresholds, actions, "
         "retention, jurisdiction and review capacity are all declared in YAML. "
         "A new jurisdiction, a new use case, or a changed risk appetite is a "
         "config edit, not a release.", 12, BODY, False, 12),
        ("An access-control breach blocks everywhere — a permission failure is not "
         "a matter of risk appetite.", 11.5, MUTED, False, 0),
    ])

    # 4 · results ---------------------------------------------------------
    s = title_slide(prs, "Results",
                    "319 labelled cases against a live judge (openai/gpt-5.6-sol). "
                    "One call in 328 fell back — 0.3%, counted and named in the report.")
    table(s, 0.61, 2.0, 12.1, [
        ["Configuration", "Catch", "False positives", "Tier 2 rate", "p95 latency", "Cost"],
        ["Rules only", "22.8%", "0.0%", "0%", "0.02 ms", "Rs 0"],
        ["Rules + classifiers", "93.5%", "7.4%", "0%", "140 ms", "Rs 0.35"],
        ["The cascade", "94.6%", "7.4%", "2.8%", "175 ms", "Rs 1.10"],
        ["A judge on every response", "92.4%", "14.8%", "100%", "5,534 ms", "Rs 30.01"],
    ], [4.0, 1.5, 2.2, 1.6, 1.6, 1.2], emphasis=3)
    textbox(s, 0.61, 3.9, 12.1, 2.4, [
        ("The cascade catches more than judging every response — 94.6% against "
         "92.4% — at 3.7% of the cost, with half the false positives and 32× lower "
         "latency.", 14, INK, True, 12),
        ("That the judge-everything baseline is worse on both is not a rhetorical "
         "win. A judge asked to re-rule on claims tier 1 already had right sometimes "
         "overrules them, and it does so in both directions. Selectivity is not only "
         "cheaper here, it is more accurate — measured against a live judge, not a "
         "stand-in.", 12, BODY, False, 0),
    ])

    # 5 · what it cannot do -----------------------------------------------
    s = title_slide(prs, "What it cannot do",
                    "Stated because a benchmark that reports only its wins is a brochure.")
    table(s, 0.61, 2.0, 12.1, [
        ["Weakness", "Rate", "Why"],
        ["Multi-hop claims", "3 of 9",
         "True only by combining two sources; entails neither alone"],
        ["Quantifier flips", "8 of 15 caught",
         '"up to X" against "at least X" — one word, the opposite meaning'],
        ["Hedged but correct", "11 of 15", "Hedging reads as distance from the source"],
        ["Two policies over their own latency budgets", "1225 / 3914 ms",
         "A judge call is seconds; a budget in hundreds of ms cannot absorb one"],
    ], [4.0, 2.1, 6.0])
    textbox(s, 0.61, 4.2, 12.1, 2.2, [
        ("The eval set is synthetic and was written by the same person who tuned the "
         "checker. An earlier version scored 100% on every case type — which is "
         "exactly why the adversarial cases exist.", 12, BODY, False, 12),
        ("These numbers are a floor on difficulty, not a ceiling on quality.", 12, INK, True, 0),
    ])

    # 6 · target users ----------------------------------------------------
    s = title_slide(prs, "Who this is for")
    table(s, 0.61, 1.9, 12.1, [
        ["Who", "What they need", "What they touch"],
        ["Platform / ML engineering", "Ship AI features without owning risk review", "A base URL change"],
        ["Risk, compliance, legal", "Evidence a decision was checked, and by what rule", "Audit export, policy YAML"],
        ["Review operations", "A queue sized to their headcount", "Escalation queue, capacity cap"],
        ["The business owner", "Cost per interaction that does not surprise them", "The cost meter"],
    ], [3.4, 5.4, 3.3])
    textbox(s, 0.61, 4.2, 12.1, 1.0, [
        ("The buyer is usually the platform team. The renewal is signed by risk.", 13, INK, True, 0),
    ])

    # 7 · the business case ----------------------------------------------
    s = title_slide(prs, "The business case",
                    "At the brief's reference volume — 30,000 interactions a week, "
                    "1.56 million a year — using measured per-response costs.")
    table(s, 0.61, 2.0, 7.2, [
        ["", "Per response", "Per year"],
        ["ControlPlane cascade", "Rs 0.0035", "Rs 5,391"],
        ["A judge on every response", "Rs 0.0941", "Rs 146,763"],
        ["Difference", "", "Rs 141,372"],
    ], [3.2, 2.0, 2.0], emphasis=1)
    textbox(s, 8.2, 2.0, 4.5, 4.3, [
        ("Verification is not the expensive part.", 13, INK, True, 4),
        ("Human review is.", 15, ACCENT, True, 12),
        ("decision_support escalates 32.3% of its traffic. At Rs 200 a review, that "
         "one use case sends 168,000 responses a year to a person: Rs 3.36 crore — "
         "against Rs 5,391 of compute for the whole estate.", 11.5, BODY, False, 10),
        ("Held to the 20% it declares it can review, the same use case costs "
         "Rs 2.08 crore. Both are real; neither is a forecast.", 11, MUTED, False, 0),
    ])
    textbox(s, 0.61, 4.25, 7.2, 2.1, [
        ("This reframes what the product is for. It is not a way to buy cheap "
         "verification — it is a way to control how much human review you have to "
         "buy, and to justify each unit of it.", 12, BODY, False, 10),
        ("Moving the escalation rate by one percentage point is worth more than the "
         "entire verification bill.", 12, INK, True, 0),
    ])
    note(s, "The eval set is 57.7% harmful by construction, so its escalation rate is an "
            "upper bound on production traffic. Cost parameters are assumptions, listed in "
            "docs/assumptions.md.")

    # 8 · break-even ------------------------------------------------------
    s = title_slide(prs, "When it pays for itself",
                    "Verification costs Rs 0.0035 per response. At the measured catch rate:")
    table(s, 0.61, 2.1, 8.0, [
        ["Use case", "Pays for itself at"],
        ["internal_copilot", "1 harmful response in 16,400"],
        ["customer_support", "1 harmful response in 109,000"],
        ["decision_support", "1 harmful response in 13.7 million"],
    ], [3.6, 4.4])
    textbox(s, 0.61, 4.0, 12.1, 2.3, [
        ("Published hallucination rates for retrieval-grounded assistants are orders "
         "of magnitude above any of these.", 13, INK, True, 12),
        ("The verification layer is not a cost decision. The review queue is.", 13, ACCENT, True, 0),
    ])

    # 9 · roadmap ---------------------------------------------------------
    s = title_slide(prs, "Roadmap")
    phases = [
        ("Phase 1 — done", "This prototype",
         "Three-tier cascade, policy layer, expected-cost router, cost meter with "
         "verified prices, hash-chained audit, OpenAI-compatible gateway, "
         "streaming-concurrent verification, feedback loop, operator dashboard, "
         "319-case benchmark with published failure modes."),
        ("Phase 2 — 1–2 quarters", "Production readiness",
         "Operator dashboard with a live escalation queue. Multi-tenant auth and "
         "RBAC. Retrieval-quality signals, so a grounded answer built on a stale "
         "chunk is distinguishable from a good one. Multi-hop grounding via "
         "evidence-set entailment. Calibration on customer data, replacing our "
         "assumed cost parameters with real ones."),
        ("Phase 3 — 2–4 quarters", "Coverage",
         "Multi-turn risk accumulation. Agent action gating — verifying tool calls, "
         "not only text. Counterfactual bias probing for disparate impact. "
         "Jurisdiction packs for EU AI Act and DPDP evidence export."),
        ("Phase 4", "Scale",
         "Distilled tier 1 models to cut the 85 ms. Regional inference for data "
         "residency. Continuous calibration from the review queue."),
    ]
    top = 1.75
    for label, name, detail in phases:
        textbox(s, 0.61, top, 2.6, 1.2, [(label, 11, ACCENT, True, 2),
                                         (name, 12.5, INK, True, 0)])
        textbox(s, 3.4, top, 9.3, 1.2, [(detail, 11, BODY, False, 0)])
        top += 1.28

    # 10 · risks ----------------------------------------------------------
    s = title_slide(prs, "Risks, and what we do about them")
    table(s, 0.61, 1.85, 12.1, [
        ["Risk", "Why it is real", "Mitigation"],
        ["Cost assumptions are ours, not a customer's",
         "Every routing decision depends on cost_of_being_wrong, and we set it",
         "One YAML value per use case; Phase 2 calibrates against real incident cost"],
        ["Alert fatigue",
         "7.4% false positives at scale is noise, and people route around noisy tools",
         "Tradeoff curve published; escalation capped by declared capacity; loop retunes"],
        ["Bad retrieval defeats grounding",
         "We check the answer against its sources, not the sources against the world",
         "Named as a limit today; retrieval-quality signals are Phase 2"],
        ["Bias detection is shallow",
         "Attribute-plus-decision heuristics cannot see disparate impact",
         "Scoped honestly; the case for Phase 3 counterfactual probing"],
        ["The judge is a model too",
         "It can be wrong, and it costs money and seconds",
         "Runs on 2.8% of traffic; verdicts logged with reasons; offline fallback"],
        ["Our own benchmark flatters us",
         "We wrote the test set and tuned against it",
         "Adversarial cases added when it scored 100%; failure modes published"],
    ], [3.5, 4.4, 4.2], size=9.5)

    # 11 · why this -------------------------------------------------------
    s = title_slide(prs, "Why this rather than the alternatives")
    textbox(s, 0.61, 2.0, 12.1, 4.3, [
        ("Runtime guardrails exist. What does not exist is one layer where all three "
         "risk dimensions are scored together, each use case runs its own risk "
         "policy, and every check is priced.", 13, BODY, False, 16),
        ("Today a team stacks an eval framework, a grounding checker, an "
         "observability platform and an inline guardrail — four products, none of "
         "which can tell you whether the check it just ran was worth what it cost.", 13, BODY, False, 20),
        ("Once verification is cheap, the real question is how much human attention "
         "to spend, and where.", 16, INK, True, 8),
        ("That is a question only a system that prices its own decisions can answer.", 14, ACCENT, True, 0),
    ])
    note(s, "Prototype, benchmark and results: github.com/m25csa003-glitch/controlplane · "
            "Every figure reproducible with python3 eval/run_eval.py")

    salutation(prs)
    prs.save(str(OUT))
    return prs


if __name__ == "__main__":
    prs = build()
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides)} slides)")
